import os
import re
import uuid
import csv
from PIL import Image
from django.conf import settings
from django.test import TestCase
from docx import Document
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Border, Side
from converter.formats import image_formats, writer_formats, table_formats, slide_formats, pdf_format

from converter.models import Conversion


class ConversionTestCase(TestCase):
    def setUp(self):
        """Initialise les variables de test et crée les dossiers nécessaires."""
        self.test_id = uuid.uuid4().hex
        self.uploads_dir = os.path.join(settings.MEDIA_ROOT, 'uploads')
        os.makedirs(self.uploads_dir, exist_ok=True)
        self.generated_files = []
        self.result_matrix = []
        self.supported_formats = image_formats + writer_formats + table_formats + slide_formats + pdf_format

    def test_conversion_matrix(self):
        """Teste toutes les combinaisons de conversion et génère une matrice de résultats."""
        for source_format in self.supported_formats:
            row = [source_format]
            for target_format in self.supported_formats:
                result = self.run_conversion_test(source_format, target_format)
                row.append(result)
            self.result_matrix.append(row)

        # Sauvegarde les résultats dans un fichier CSV
        self.save_result_matrix()

    def run_conversion_test(self, source_format, target_format):
        """
        Effectue un test de conversion entre un format source et cible.

        Args:
            source_format (str): Format source.
            target_format (str): Format cible.

        Returns:
            str: Résultat de la conversion ("S", "N", "E" ou "#").
        """
        if source_format == target_format:
            return "#"

        file_name = f"test_{source_format}_{self.test_id}.{source_format}"
        file_path = self.create_fake_file(file_name, source_format)
        conversion = self.create_conversion(file_path, source_format, target_format)

        try:
            conversion.convert_file()
            if conversion.converted:
                self.generated_files.append(conversion.output_file.path)
                print(f"🟩 Conversion réussie: {source_format} -> {target_format}")
                return "S" # Supported
            else:
                error_message = conversion.error_message or "Unknown error"
                if re.search("Unsupported conversion", error_message):
                    print(f"🟨 Conversion non supportée: {source_format} -> {target_format}")
                    return "N" # Not supported
                else:
                    print(f"🟥 Erreur lors de la conversion {source_format} -> {target_format}: {error_message}")
                    return f"E" # Error
        except Exception as e:
            print(f"🟥 Exception lors de la conversion {source_format} -> {target_format}: {str(e)}")
            return f"E" # Error

    def create_conversion(self, file_path, source_format, target_format):
        """Crée une instance de Conversion."""
        self.generated_files.append(file_path)
        return Conversion(input_file=file_path, source_format=source_format, target_format=target_format)

    def create_fake_file(self, file_name, file_type):
        """Crée un fichier factice pour un format donné."""
        file_path = os.path.join(self.uploads_dir, file_name)

        if file_type in image_formats:
            image = Image.new('RGB', (100, 100), color=(255, 0, 0))
            image.save(file_path)
            
        elif file_type in writer_formats:
            if file_type == "docx":
                doc = Document()
                doc.add_paragraph("This is a test file.")
                doc.save(file_path)
            else:
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write("This is a test file.")
        
        elif file_type in table_formats:
            if file_type == "xlsx":
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws["A1"] = "This is a test file."
                wb.save(file_path)
            else:
                pass # TODO: Add support for other table formats
           
        elif file_type in slide_formats:
            if file_type == "pptx":
                from pptx import Presentation
                presentation = Presentation()
                slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                slide.shapes.title.text = "Test Presentation"
                slide.placeholders[1].text = "This is a test file."
                presentation.save(file_path)  
            else:
                pass # TODO: Add support for other slide formats
            
        elif file_type == pdf_format:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.cell(200, 10, txt="This is a test file.", ln=True)
            pdf.output(file_path)
        
        self.generated_files.append(file_path)
        return file_path

    def save_result_matrix(self):
        """Sauvegarde la matrice des résultats dans un fichier XLSX avec des couleurs."""
        result_path = os.path.join(settings.MEDIA_ROOT, "conversions.xlsx")
        
        # Créer un nouveau classeur Excel
        wb = Workbook()
        ws = wb.active
        ws.title = "Conversion Results"

        # Définir les couleurs
        color_map = {
            "S": "00FF00",  # Vert pour Supported
            "N": "FFFF00",  # Jaune pour Not Supported
            "E": "FF0000",  # Rouge pour Error
            "#": "CCCCCC",  # Gris pour Identique
        }

        # Définir le style de bordure
        thin_top_border = Border(top=Side(style="thin"))
    
        # En-tête des colonnes
        ws.append(["Source \\ Target"] + self.supported_formats)

        # Appliquer la bordure à l'en-tête
        for col_num in range(1, len(self.supported_formats) + 2):  # +2 pour inclure la colonne Source \ Target
            cell = ws.cell(row=1, column=col_num)
            cell.border = thin_top_border

        # Remplir les données avec couleurs et bordures
        for i, row in enumerate(self.result_matrix, start=2):  # Commence à la 2e ligne (1 = en-têtes)
            for j, cell_value in enumerate(row, start=1):      # Commence à la 1ère colonne
                cell = ws.cell(row=i, column=j, value=cell_value)

                # Appliquer la couleur si le résultat est dans color_map
                if cell_value in color_map:
                    cell.fill = PatternFill(start_color=color_map[cell_value], end_color=color_map[cell_value], fill_type="solid")
                
                # Appliquer une bordure supérieure
                cell.border = thin_top_border
                
        # Enregistrer le fichier XLSX
        wb.save(result_path)
        print(f"🟦 Résultats des conversions sauvegardés dans {result_path}")

    def tearDown(self):
        """Nettoie les fichiers générés après les tests."""
        for file_path in self.generated_files:
            if os.path.exists(file_path):
                os.remove(file_path)
